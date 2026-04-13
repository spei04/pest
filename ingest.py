"""
Ingests documents from sources/ into a ChromaDB vector store.
Supported formats: .docx, .pptx, .pdf, .mp4

Run once (or re-run to update):
    python ingest.py

Re-running is safe — upsert with deterministic IDs means no duplicates.
"""

import hashlib
import os
import subprocess
import sys
import tempfile
from pathlib import Path

import chromadb
import pdfplumber
from docx import Document
from dotenv import load_dotenv
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from transformers import logging as hf_logging
hf_logging.set_verbosity_error()

load_dotenv()

SOURCES_DIR = Path("sources")
CHROMA_DIR = os.environ.get("CHROMA_DIR", "./chroma_db")
COLLECTION_NAME = "pest_docs"
EMBED_MODEL = "all-MiniLM-L6-v2"  # free, local, no API key needed
EMBED_BATCH_SIZE = 100
CHUNK_SIZE = 1500
CHUNK_OVERLAP = 200


# ---------------------------------------------------------------------------
# Parsers
# ---------------------------------------------------------------------------

def parse_docx(path: Path) -> str:
    doc = Document(path)
    return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def parse_pdf(path: Path) -> str:
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text and text.strip():
                pages.append(text.strip())
    return "\n\n".join(pages)


def parse_pptx(path: Path) -> list[str]:
    """Returns one chunk per slide (slides are natural semantic units)."""
    prs = Presentation(path)
    slide_texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = "\n".join(parts)
        if len(text) >= 20:  # skip near-empty / image-only slides
            slide_texts.append(text)
    return slide_texts


def parse_mp4(path: Path) -> str:
    """Extracts audio via ffmpeg, transcribes with Whisper."""
    import whisper

    model = whisper.load_model("small")

    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
        tmp_path = tmp.name

    try:
        print(f"  Extracting audio from {path.name} (this may take a moment)...")
        subprocess.run(
            ["ffmpeg", "-y", "-i", str(path), "-ar", "16000", "-ac", "1", tmp_path],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        print(f"  Transcribing with Whisper small (may take several minutes for long recordings)...")
        result = model.transcribe(tmp_path)
        return result["text"]
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_text(text: str) -> list[str]:
    """
    Split on paragraph boundaries first; hard-split any paragraph that
    exceeds CHUNK_SIZE with CHUNK_OVERLAP.
    """
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks = []
    current = ""

    for para in paragraphs:
        # If adding this paragraph keeps us under the limit, accumulate it
        if len(current) + len(para) + 2 <= CHUNK_SIZE:
            current = (current + "\n\n" + para).strip() if current else para
        else:
            if current:
                chunks.append(current)
            # If the paragraph itself is larger than CHUNK_SIZE, hard-split it
            if len(para) > CHUNK_SIZE:
                start = 0
                while start < len(para):
                    chunks.append(para[start : start + CHUNK_SIZE])
                    start += CHUNK_SIZE - CHUNK_OVERLAP
                current = ""
            else:
                current = para

    if current:
        chunks.append(current)

    return chunks


# ---------------------------------------------------------------------------
# Embedding
# ---------------------------------------------------------------------------

def embed_texts(texts: list[str], model: SentenceTransformer) -> list[list[float]]:
    return model.encode(texts, batch_size=EMBED_BATCH_SIZE, show_progress_bar=False).tolist()


# ---------------------------------------------------------------------------
# Chunk ID
# ---------------------------------------------------------------------------

def chunk_id(source_rel: str, index: int) -> str:
    raw = f"{source_rel}::chunk_{index}"
    return hashlib.md5(raw.encode()).hexdigest()


# ---------------------------------------------------------------------------
# Ingest a single file
# ---------------------------------------------------------------------------

def ingest_file(path: Path, collection, model: SentenceTransformer) -> int:
    ext = path.suffix.lower()
    source_rel = str(path.relative_to(SOURCES_DIR))

    print(f"Processing: {source_rel}")

    try:
        if ext == ".docx":
            text = parse_docx(path)
            chunks = chunk_text(text)
        elif ext == ".pdf":
            text = parse_pdf(path)
            chunks = chunk_text(text)
        elif ext == ".pptx":
            chunks = parse_pptx(path)
        elif ext == ".mp4":
            text = parse_mp4(path)
            chunks = chunk_text(text)
        else:
            print(f"  Skipping unsupported format: {ext}")
            return 0
    except Exception as e:
        print(f"  ERROR: {e}")
        return 0

    if not chunks:
        print(f"  Warning: no text extracted, skipping.")
        return 0

    embeddings = embed_texts(chunks, model)

    ids = [chunk_id(source_rel, i) for i in range(len(chunks))]
    metadatas = [
        {"source": source_rel, "file_type": ext.lstrip("."), "chunk_index": i}
        for i in range(len(chunks))
    ]

    collection.upsert(
        ids=ids,
        embeddings=embeddings,
        documents=chunks,
        metadatas=metadatas,
    )

    print(f"  Ingested {len(chunks)} chunks.")
    return len(chunks)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    if not SOURCES_DIR.exists():
        print(f"ERROR: '{SOURCES_DIR}' directory not found.")
        sys.exit(1)

    supported = {".docx", ".pptx", ".pdf", ".mp4"}
    files = sorted(p for p in SOURCES_DIR.rglob("*") if p.is_file() and p.suffix.lower() in supported)

    if not files:
        print("No supported files found in sources/")
        sys.exit(1)

    print(f"Found {len(files)} files.\n")

    print(f"Loading embedding model '{EMBED_MODEL}'...")
    embed_model = SentenceTransformer(EMBED_MODEL)

    chroma = chromadb.PersistentClient(path=CHROMA_DIR)
    collection = chroma.get_or_create_collection(
        name=COLLECTION_NAME,
        metadata={"hnsw:space": "cosine"},
    )

    total = 0
    for f in files:
        total += ingest_file(f, collection, embed_model)

    print(f"\nDone. {total} total chunks stored in '{CHROMA_DIR}'.")


if __name__ == "__main__":
    main()
